"""DarkLab Media Generator: produce Word docs and PPTX from synthesized research.

Invoked by OpenClaw node-host via:
  python3 -m leader.media_gen '{"task_type":"media_gen","payload":{...}}'

Uses python-docx for Word documents and python-pptx for presentations.
"""
from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path

import structlog

from shared.models import Task, TaskResult
from shared.config import settings
from shared.audit import log_event
from shared.node_bridge import run_agent

logger = structlog.get_logger("darklab.media_gen")


def _generate_docx(synthesis: dict, output_dir: Path, task_id: str) -> Path:
    """Generate a Word document from synthesis data."""
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    title_para = doc.add_heading("DarkLab Research Report", level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    doc.add_paragraph("")

    # Executive Summary
    if synthesis.get("executive_summary"):
        doc.add_heading("Executive Summary", level=1)
        doc.add_paragraph(synthesis["executive_summary"])

    # Key Findings
    findings = synthesis.get("key_findings", [])
    if findings:
        doc.add_heading("Key Findings", level=1)
        for finding in findings:
            doc.add_paragraph(finding, style="List Bullet")

    # Methodology Validation
    if synthesis.get("methodology_validation"):
        doc.add_heading("Methodology Validation", level=1)
        doc.add_paragraph(synthesis["methodology_validation"])

    # Data Consistency
    consistency = synthesis.get("data_consistency", {})
    if consistency:
        doc.add_heading("Data Consistency", level=1)
        score = consistency.get("score", "N/A")
        doc.add_paragraph(f"Consistency Score: {score}")
        for issue in consistency.get("issues", []):
            doc.add_paragraph(issue, style="List Bullet")

    # Recommendations
    recommendations = synthesis.get("recommendations", [])
    if recommendations:
        doc.add_heading("Recommendations", level=1)
        for i, rec in enumerate(recommendations, 1):
            doc.add_paragraph(f"{i}. {rec}")

    # Full Narrative
    if synthesis.get("full_narrative"):
        doc.add_heading("Detailed Analysis", level=1)
        for paragraph_text in synthesis["full_narrative"].split("\n\n"):
            if paragraph_text.strip():
                doc.add_paragraph(paragraph_text.strip())

    output_path = output_dir / f"report_{task_id}.docx"
    doc.save(str(output_path))
    return output_path


def _generate_pptx(synthesis: dict, output_dir: Path, task_id: str) -> Path:
    """Generate a PowerPoint presentation from synthesis data."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "DarkLab Research Report"
    slide.placeholders[1].text = datetime.now().strftime("%Y-%m-%d")

    # Executive Summary slide
    if synthesis.get("executive_summary"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary"
        slide.placeholders[1].text = synthesis["executive_summary"][:800]

    # Key Findings slide
    findings = synthesis.get("key_findings", [])
    if findings:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Key Findings"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = findings[0] if findings else ""
        for finding in findings[1:]:
            p = tf.add_paragraph()
            p.text = finding
            p.level = 0

    # Recommendations slide
    recommendations = synthesis.get("recommendations", [])
    if recommendations:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Recommendations"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = recommendations[0] if recommendations else ""
        for rec in recommendations[1:]:
            p = tf.add_paragraph()
            p.text = rec
            p.level = 0

    # Data Consistency slide
    consistency = synthesis.get("data_consistency", {})
    if consistency and consistency.get("score") is not None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Data Consistency"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = f"Consistency Score: {consistency.get('score', 'N/A')}"
        for issue in consistency.get("issues", []):
            p = tf.add_paragraph()
            p.text = issue
            p.level = 1

    output_path = output_dir / f"slides_{task_id}.pptx"
    prs.save(str(output_path))
    return output_path


async def handle(task: Task) -> TaskResult:
    synthesis = task.payload.get("synthesis", {})
    output_types = task.payload.get("output_types", ["report"])

    if not synthesis:
        return TaskResult(
            task_id=task.task_id,
            agent_name="MediaGenAgent",
            status="error",
            result={"error": "No synthesis data provided. Include 'synthesis' in payload."},
        )

    output_dir = settings.artifacts_dir / f"media_{task.task_id}"
    output_dir.mkdir(parents=True, exist_ok=True)

    deliverables = []
    artifacts = []

    if "report" in output_types:
        try:
            docx_path = _generate_docx(synthesis, output_dir, task.task_id)
            deliverables.append({"type": "report", "path": str(docx_path), "format": "docx"})
            artifacts.append(str(docx_path))
            logger.info(f"Generated Word report: {docx_path}")
        except Exception as e:
            logger.error(f"DOCX generation failed: {e}")
            deliverables.append({"type": "report", "error": str(e)})

    if "presentation" in output_types:
        try:
            pptx_path = _generate_pptx(synthesis, output_dir, task.task_id)
            deliverables.append({"type": "presentation", "path": str(pptx_path), "format": "pptx"})
            artifacts.append(str(pptx_path))
            logger.info(f"Generated PPTX: {pptx_path}")
        except Exception as e:
            logger.error(f"PPTX generation failed: {e}")
            deliverables.append({"type": "presentation", "error": str(e)})

    log_event("media_gen", task_id=task.task_id, types=output_types, count=len(artifacts))

    return TaskResult(
        task_id=task.task_id,
        agent_name="MediaGenAgent",
        status="ok" if artifacts else "error",
        result={"deliverables": deliverables},
        artifacts=artifacts,
    )


if __name__ == "__main__":
    run_agent(handle, agent_name="MediaGenAgent")
